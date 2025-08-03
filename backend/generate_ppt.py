import json
import sys
import os
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
import google.generativeai as genai
from dotenv import load_dotenv

# Load .env environment variables
load_dotenv()
sys.stdout.reconfigure(encoding='utf-8') # Prevent encoding issues in console

# === Function to create PPT from JSON ===
def generate_ppt_from_template(json_content, template_path, output_path):
    print(f"üî∏ Checking template path: {template_path}")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"‚ùå Template file not found at: {template_path}")
    print("üî∏ Loading PowerPoint template...")
    prs = Presentation(template_path)
    try:
        slides_data = json.loads(json_content)
    except json.JSONDecodeError as e:
        raise ValueError(f"‚ùå Invalid JSON passed to generate_ppt_from_template: {e}")

    print("üî∏ Populating slides...")
    for i, slide_data in enumerate(slides_data.get("slides", [])):
        # Use different layouts for title slide and content slides
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

        if "content" in slide_data:
            try:
                # Assuming placeholder[1] is the content placeholder
                if len(slide.shapes.placeholders) > 1:
                    content_box = slide.shapes.placeholders[1].text_frame
                    content_box.clear()
                    for point in slide_data["content"]:
                        content_box.add_paragraph().text = point.strip()
                else:
                    print(f"‚ö†Ô∏è Slide {i+1} has no content placeholder at index 1. Skipping content.")
            except Exception as e:
                print(f"‚ö†Ô∏è Could not add bullet points to slide {i+1}: {e}")
    prs.save(output_path)
    print(f"‚úÖ Presentation saved at: {output_path}")

# === Main Execution ===
if __name__ == "__main__":
    if len(sys.argv) < 3:
        raise ValueError("‚ùå Missing arguments. Usage: python generate_ppt.py <input_pdf_path> <output_pptx_path>")

    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2] # Get output path from command line argument

    if not os.path.exists(input_pdf):
        raise FileNotFoundError(f"‚ùå PDF not found at: {input_pdf}")

    print("üîπ Extracting text from PDF...")
    loader = PyPDFLoader(input_pdf)
    text = " ".join([doc.page_content for doc in loader.load()])[:8000] # Limit input

    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

    # First model: get summary
    print("üîπ Generating Summary...")
    summary_model = genai.GenerativeModel("gemini-2.0-flash-exp", generation_config={"response_mime_type": "text/plain"})
    summary_chat = summary_model.start_chat()
    summary_response = summary_chat.send_message(text).text.strip()

    if not summary_response:
        raise ValueError("‚ùå Empty AI summary response")
    print("üîπ AI Summary Response:")
    print(summary_response)

    # Second model: convert summary to JSON
    print("üîπ Generating JSON...")
    json_model = genai.GenerativeModel(
        model_name="gemini-2.0-flash-exp",
        generation_config={"response_mime_type": "application/json"},
        system_instruction="""Convert the research paper summary into a structured JSON presentation.
        Each slide should contain a title and up to 4-5 bullet points.
        Output format:
        {
            "presentation_title": "Research Paper Title",
            "slides": [
                {"type": "title", "title": "Title Slide", "subtitle": "Subtitle"},
                {"type": "content", "title": "Slide Title", "content": ["Point 1", "Point 2"]}
            ]
        }
        """
    )
    json_chat = json_model.start_chat()
    try:
        json_text = json_chat.send_message(summary_response).text.strip()
    except Exception as e:
        raise RuntimeError(f"‚ùå Gemini JSON generation failed: {e}")

    if not json_text:
        raise ValueError("‚ùå Empty AI JSON response")
    print("üîπ AI JSON Response:")
    print(json_text)

    # Validate JSON
    print("üîπ Validating JSON...")
    try:
        json.loads(json_text)
    except json.JSONDecodeError as e:
        print("‚ùå Bad JSON content:")
        print(json_text)
        raise ValueError(f"‚ùå Invalid JSON from Gemini: {e}")

    # Prepare paths
    # Assuming Template3-ppt.pptx is in the same directory as the script
    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Template3-ppt.pptx")
    
    print("üîπ Generating PowerPoint...")
    generate_ppt_from_template(json_text, template_path, output_pptx) # Pass the output_pptx path

    print(f"\n‚úÖ PowerPoint generated successfully: {output_pptx}")
