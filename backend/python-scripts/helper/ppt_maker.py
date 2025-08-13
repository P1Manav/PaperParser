import json
import os
import re
from dotenv import load_dotenv
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from PIL import Image



def make_ppt_from_data(template_path:str , output_path:str):

   
    pre = Presentation(template_path)
    with open(r"final_ppt_data.json", "r", encoding="utf-8") as f:
        slides_data = json.load(f)
    # print(slides_data)

    # for Title Slide of the ppt

    title_slide_data = slides_data[0]
    layout_index = title_slide_data["layout_index"]
    placeholder = title_slide_data["placeholders"]


    slide_layout = pre.slide_layouts[layout_index]
    slide = pre.slides.add_slide(slide_layout)


    if "title" in placeholder:
        try:
            title_shape = slide.placeholders[placeholder["title"]]
            title_shape.text = title_slide_data["slide_title"]

            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(32)

        except:

            print(f"Title page main heading not found")

    if "content" in placeholder:
        try:
            content_shape = slide.placeholders[placeholder["content"]]
            content_shape.text_frame.text =""

            for bullet in title_slide_data["bullet_points"]:
                p = content_shape.text_frame.add_paragraph()
                p.text = bullet
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(14)
        except:
            print(f"Title page sub heading not found")


    for slide_data in slides_data[1:]:

        layout_index = slide_data["layout_index"]
        placeholders = slide_data["placeholders"]

        slide_layout = pre.slide_layouts[layout_index]
        slide = pre.slides.add_slide(slide_layout)
        if "title" in placeholders:
            try:
                title_shape = slide.placeholders[placeholders["title"]]
                title_shape.text = slide_data["slide_title"]

                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(30)
                        run.font.bold = True

            except IndexError:
                print(f"Title placeholder index {placeholders['title']} not found")

        if "content" in placeholders:
            try:
                content_shape = slide.placeholders[placeholders["content"]]

                content_shape.text_frame.text = ""

     

                raw_bullets = slide_data.get("bullet_points", [])
                cleaned_bullets = [bullet.replace("**", "").strip() for bullet in raw_bullets]


                # Add bullets to the content placeholder
                for bullet in cleaned_bullets:
                    p = content_shape.text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(14)
                        run.font.bold = True

                ph = slide.placeholders[placeholders["content"]]





            except IndexError:
                print(f"Content placeholder index {placeholders['content']} not found")

        if "image" in placeholders and slide_data["image_path"] != "null":
            try:

                ph = slide.placeholders[placeholders["image"]]
                img_path = "images/" + slide_data["image_path"]
                # getting ph dimension
                ph_left = ph.left
                ph_top = ph.top
                ph_width = ph.width
                ph_height = ph.height

                # original image size
                im = Image.open(img_path)
                img_width, img_height = im.size
                img_ratio = img_width / img_height
                ph_ratio = ph_width / ph_height

                if img_ratio > ph_ratio:  # original image is wider
                    new_width = ph_width
                    new_height = ph_width / img_ratio
                else:
                    new_height = ph_height
                    new_width = ph_height * img_ratio
                # scale_w = 1.2  # 5% bigger
                # scale_h = 1.2
                # new_width = int(new_width * scale_w)
                # new_height = int(new_height * scale_h)
                left = ph_left + (ph_width - new_width) // 2
                top = ph_top + (ph_height - new_height) // 2
                slide.shapes.add_picture(
                    img_path, left, top, width=new_width, height=new_height
                )

            except Exception as e:
                print(f"Image not inserted on slide '{slide_data['slide_title']}': {e}")


    pre.save(output_path)
    print(f"Presentation saved at: {output_path}")




