# import json
# import sys
# import os
# from pptx import Presentation
# from langchain_community.document_loaders import PyPDFLoader
# import google.generativeai as genai
# from dotenv import load_dotenv

# # Load .env environment variables
# load_dotenv()
# sys.stdout.reconfigure(encoding='utf-8') # Prevent encoding issues in console
# genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# # === Function to create PPT from JSON ===
# def generate_ppt_from_template(json_content, template_path, output_path):
#     print(f"🔸 Checking template path: {template_path}")
#     if not os.path.exists(template_path):
#         raise FileNotFoundError(f"❌ Template file not found at: {template_path}")
#     print("🔸 Loading PowerPoint template...")
#     prs = Presentation(template_path)
#     try:
#         slides_data = json.loads(json_content)
#     except json.JSONDecodeError as e:
#         raise ValueError(f"❌ Invalid JSON passed to generate_ppt_from_template: {e}")

#     print("🔸 Populating slides...")
#     for i, slide_data in enumerate(slides_data.get("slides", [])):
#         # Use different layouts for title slide and content slides
#         layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
#         slide = prs.slides.add_slide(layout)

#         if slide.shapes.title:
#             slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

#         if "content" in slide_data:
#             try:
#                 # Assuming placeholder[1] is the content placeholder
#                 if len(slide.shapes.placeholders) > 1:
#                     content_box = slide.shapes.placeholders[1].text_frame
#                     content_box.clear()
#                     for point in slide_data["content"]:
#                         content_box.add_paragraph().text = point.strip()
#                 else:
#                     print(f"⚠️ Slide {i+1} has no content placeholder at index 1. Skipping content.")
#             except Exception as e:
#                 print(f"⚠️ Could not add bullet points to slide {i+1}: {e}")
#     prs.save(output_path)
#     print(f"✅ Presentation saved at: {output_path}")

# # === Main Execution ===
# if __name__ == "__main__":
#     if len(sys.argv) < 3:
#         raise ValueError("❌ Missing arguments. Usage: python generate_ppt.py <input_pdf_path> <output_pptx_path>")

#     input_pdf = sys.argv[1]
#     output_pptx = sys.argv[2] # Get output path from command line argument
#     print(input_pdf)
#     print(output_pptx)


#     if not os.path.exists(input_pdf):
#         raise FileNotFoundError(f"❌ PDF not found at: {input_pdf}")

#     print("🔹 Extracting text from PDF...")
#     loader = PyPDFLoader(input_pdf)
#     text = " ".join([doc.page_content for doc in loader.load()])[:8000] # Limit input

#     genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

#     # First model: get summary
#     print("🔹 Generating Summary...")
#     summary_model = genai.GenerativeModel("gemini-2.0-flash-exp", generation_config={"response_mime_type": "text/plain"})
#     summary_chat = summary_model.start_chat()
#     summary_response = summary_chat.send_message(text).text.strip()

#     if not summary_response:
#         raise ValueError("❌ Empty AI summary response")
#     print("🔹 AI Summary Response:")
#     print(summary_response)

#     # Second model: convert summary to JSON
#     print("🔹 Generating JSON...")
#     json_model = genai.GenerativeModel(
#         model_name="gemini-2.0-flash-exp",
#         generation_config={"response_mime_type": "application/json"},
#         system_instruction="""Convert the research paper summary into a structured JSON presentation.
#         Each slide should contain a title and up to 4-5 bullet points.
#         Output format:
#         {
#             "presentation_title": "Research Paper Title",
#             "slides": [
#                 {"type": "title", "title": "Title Slide", "subtitle": "Subtitle"},
#                 {"type": "content", "title": "Slide Title", "content": ["Point 1", "Point 2"]}
#             ]
#         }
#         """
#     )
#     json_chat = json_model.start_chat()
#     try:
#         json_text = json_chat.send_message(summary_response).text.strip()
#     except Exception as e:
#         raise RuntimeError(f"❌ Gemini JSON generation failed: {e}")

#     if not json_text:
#         raise ValueError("❌ Empty AI JSON response")
#     print("🔹 AI JSON Response:")
#     print(json_text)

#     # Validate JSON
#     print("🔹 Validating JSON...")
#     try:
#         json.loads(json_text)
#     except json.JSONDecodeError as e:
#         print("❌ Bad JSON content:")
#         print(json_text)
#         raise ValueError(f"❌ Invalid JSON from Gemini: {e}")

#     # Prepare paths
#     # Assuming Template3-ppt.pptx is in the same directory as the script
#     template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Template3-ppt.pptx")
    
#     print("🔹 Generating PowerPoint...")
#     generate_ppt_from_template(json_text, template_path, output_pptx) # Pass the output_pptx path

#     print(f"\n✅ PowerPoint generated successfully: {output_pptx}")






# /////////////////////////////////////////////////////////


import json
import os
import re
import sys
from pptx import Presentation
from langchain_community.document_loaders import PyMuPDFLoader
from google import genai
from dotenv import load_dotenv
from helper.image_extractor import extract_combined_images_with_captions
from helper.ppt_maker import make_ppt_from_data

load_dotenv()



# -------------------


# -----------------------





#  Phase -1 Data extraction from the ppt //////////////////////////////////////////////////

def generate_ppt(input_path,template_path, output_path):

    print(f"Get the input {input_path}\n")
    print(f"Get the template {template_path}\n")
    print(f"Get the output {template_path}\n")
    


    file = input_path
    loader = PyMuPDFLoader(file)
    docs = loader.load()
    text = ""
    for doc in docs:
        text += doc.page_content

    print("\nPhase-1: Text Extraction is completed\n")
    # Phase -2 Take the extracted data and generate a ppt text for according to pages /////////////////////////////////////////////////////////////////////////

    # type_pdf = input("\n.....Enter the pdf type... \n" )
 
    # input of pdf type i.e research paper or any standard doc


    # length_of_ppt = input("\n..Enter the lenght of ppt means short / medium / long..\n")

    length_of_ppt = "long"
    system_prompt_1 = f"""


    You are a smart document structuring assistant. Your task is to take the complete extracted text from a PDF file and preprocess it in a way that it becomes well-organized, logically grouped, and ready for PowerPoint slide generation in the next phase.

    Here are the inputs:

    -----------------------------------
    TYPE OF PDF: Research Paper
    TARGET PPT LENGTH: {length_of_ppt} (short = 5-6 slides, medium = 7-8 slides, long = 10-12 slides)

    FULL RAW TEXT EXTRACTED FROM PDF:
    {text}
    -----------------------------------

    Your goals:

    1. Understand and analyze the structure of the document based on its type:
    - For research papers, identify and label sections like: Abstract, Introduction, Problem Statement, Literature Review, Methodology, Results, Conclusion, Future Work, etc.
    - For standard documents (e.g. proposals, business plans), group and label based on logical themes or topics like Overview, Objectives, Key Concepts, Strategies, Implementation, Case Studies, etc.
    - If type is unknown or custom, use best judgment to extract and organize major topics or logical chunks.

    2. Group the content into structured labeled sections:
    - Use headers like `## Section: Introduction`, `## Section: Results`, etc.
    - Each section should contain the full detailed text relevant to that label.
    - Preserve important technical or formal tone (especially for research-type PDFs)

    3. Adjust the depth and number of sections based on the requested presentation length:
    - If short, limit to 5-6 major sections with concise, focused content.
    - If medium, allow 7-8 moderately detailed sections.
    - If long, include 10-12+ granular sections and sub-sections.

    4. Preprocess the content:
    - Clean up repetitions or OCR errors.
    - Merge fragmented lines/sentences.
    - Retain any bullet points, formulas, or figure captions as-is.
    - If figure captions are found, label them clearly as `[Figure Caption: ...]`

    5. Do NOT generate slides yet. Just return a well-organized, clean, and structured text document — ready for slide conversion in the next phase.

    Return the output in clear Markdown or readable labeled format.

    Now process the content accordingly.



    """

    gemini_api_key = os.getenv("gemini_api_key")
    client = genai.Client(api_key=gemini_api_key)

    response = client.models.generate_content(
        model="gemini-2.5-pro",
        contents=system_prompt_1,
    )

    # print(response.text)


    pre_process_text = response.text
    print("\nPhase-2: Text Processing Completed \n")


    # Phase -3 from the pdf find the images /////////////////////////////


    extract_combined_images_with_captions(file)

    print("\nPhase-3: Retreive the Images from the Given PDF\n")


    # Phase -4 from the text find all the figure_captions ////////////////////////////////////


    system_prompt_2 = f"""


    You are a figure caption extractor for research papers.

    Below is the preprocessed and cleaned full text of a research paper. Your task is to accurately extract all figure captions, which follow standard academic formatting such as:

    - "Figure 1: ..."  
    - "Fig. 2: ..."  
    - Sometimes long captions span multiple lines after the label.
    INPUT TEXT:
    {pre_process_text}

    Your goals:

    1. Extract every figure caption that starts with:
    - "Figure X:" (where X is a number)
    - OR "Fig. X:" (alternative abbreviation)

    2. If a caption spans multiple lines, combine all the lines until the next paragraph or until the next "Figure"/"Fig." starts.

    3. Return a clean ordered list of all figure captions found, in the format:
    - Caption Number (e.g., Figure 4)
    - Full Caption Text (include all lines that belong to it)

    4. Do not modify or summarize captions. Just extract exactly what is written after the "Figure" or "Fig." label.

    5. Skip any unrelated content.

    Format your output like this:
    Figure 1
    Full caption text...

    Figure 2
    Another caption text.

    """


    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=system_prompt_2,
    )

    figure_captions = response.text

    # print(figure_captions)
    print("\nPhase-4: All Figures Captions are extracted \n")


    # phase 5 /////////////////////////////////////////


    with open(r"images\image_captions.json", "r", encoding="utf-8") as f:
        figure_captions_data = json.load(f)


    system_prompt_3 = f"""
    You are a smart assistant helping to clean and complete image captions for a research paper presentation.

    ### INPUTS:

    1. A dictionary named `image_captions`:
    - Each key is an image filename (e.g., "image23.png").
    - Each value is a partial, incomplete, or noisy caption string extracted from the figure.

    2. A list named `figure_captions`:
    - This contains all the clean, complete, original figure captions from the research paper.

    ### TASK:

    Your job is to intelligently match and update the partial captions using the full captions.

    ### INSTRUCTIONS:

    - For each entry in `image_captions`, try to find a matching full caption from `figure_captions`.
    - If a confident match is found (based on figure number, meaningful phrase overlap, etc.), replace the partial caption with the exact matched full caption.
    - Do not modify or include:
    - Entries where the original value is "No Figure caption found"
    - Entries where a confident match is not found in `figure_captions`

    ### OUTPUT REQUIREMENTS:

    - Return a dictionary in strict JSON format
    - Only include entries that were confidently matched and corrected
    - All keys (image filenames) must be unique
    - All values (captions) must be unique and directly from `figure_captions`
    - Do not generate or modify captions on your own
    - Ensure the final output is strictly parseable as JSON

    ### DATA:

    Partial image captions (dictionary):
    {figure_captions_data}

    Full figure captions (list):
    {figure_captions}

    ### STRICT OUTPUT FORMAT

    Return only the corrected dictionary as valid JSON:

    {{
        "image1.png": "Figure 1: Full caption here.",
        "image2.png": "Figure 2: Another caption."
    }}

    IMPORTANT:
    - Output only the corrected dictionary in valid JSON format.
    - Do not include any explanation, code block formatting, or markdown like ```json.
    - Return only raw JSON.
    """


    response = client.models.generate_content(
        model="gemini-2.5-pro", contents=system_prompt_3
    )

    image_captions_text = response.text.strip()

    print("\nPhase-5: Preprocess the figure captions and correct it \n")


    if image_captions_text.startswith("```json"):
        image_captions_text = re.sub(
            r"^```json\s*|\s*```$", "", image_captions_text.strip(), flags=re.MULTILINE
        )
    # print(image_captions_text)

    try:
        image_caption_dict = json.loads(image_captions_text)
    except:
        print("Not valid json format")
        image_caption_dict = {}


    with open("image_captions.json", "w", encoding="utf-8") as f:
        json.dump(image_caption_dict, f, indent=4, ensure_ascii=False)


    # Phase 6 Actual generation of the ppt data/////////////////////////////////////////////////////////

    system_prompt_4 = f"""

    You are a highly skilled research assistant and expert PowerPoint slide designer. Your task is to generate a structured JSON array representing a **professional, visually pleasing PowerPoint presentation** for a research paper. You will use the following two inputs


    1. **preprocess_text**: A clean, linearized version of the full research paper. It contains all major sections — Abstract, Introduction, Methodology, Experiments, Results, Discussion, Figures,etc  and Conclusion — in plain text format.Below is complete pre_process_text

    {pre_process_text}

    2. **image_caption_dict**: A dictionary where the keys are image filenames (e.g., "image4.png") and the values are the corresponding cleaned figure captions. These images should only be placed on slides **where appropriate and contextually meaningful**.Given Below

    {image_caption_dict}

    YOUR OBJECTIVE

    Generate a **parallel, clean, and consistent slide deck** structure in the form of a JSON array. Each object in the array represents a slide, including:

    - Title
    - Bullet points (if textual)
    - Image (only when relevant)
    - Professional styling attributes (font size, color)

    The goal is to create a beautiful, reader-friendly research presentation that summarizes the key contributions of the paper section by section.


    STRICT SLIDE STRUCTURE FORMAT

    {{
    "slide_title": "Introduction",

    "content_type": "bullet_points",  // or "image" or "mixed"
    "bullet_points": [
        "Each bullet should be concise and clear.",
        "Focus on key findings or motivations.",
        "Use simple language without jargon."
    ],

    "image_path": null  // or e.g., "image4.png" only if the figure caption is highly relevant
    }}


    """


    response = client.models.generate_content(
        model="gemini-2.5-pro", contents=system_prompt_4
    )

    ppt_json_text = response.text.strip()

    # print(ppt_json_text)

    print("\nPhase-6: PPT data is generated\n")

    if ppt_json_text.startswith("```json"):
        ppt_json_text = re.sub(
            r"^```json\s*|\s*```$", "", ppt_json_text.strip(), flags=re.MULTILINE
        )

    try:
        ppt_json_dict = json.loads(ppt_json_text)
    except:
        print("Not valid json format")
        ppt_json_dict = {}


    with open("ppt_data.json","w",encoding="utf-8") as f:
        json.dump(ppt_json_dict,f ,indent=4,ensure_ascii=False)





    #Phase 7 Alignment of the ppt 

    pre = Presentation(template_path)
    ppt_template_layout_info = ""
    for i, layout in enumerate(pre.slide_layouts):
        ppt_template_layout_info += f"Layout {i} Name: {layout.name}\n"
        for j, placeholder in enumerate(layout.placeholders):
            ppt_template_layout_info += f"Placeholder {j} - idx: {placeholder.placeholder_format.idx}, type: {placeholder.placeholder_format.type}\n"


    system_prompt_5 = f"""

    You're an AI assistant helping to generate a PowerPoint JSON structure. Use the provided dictionary (slide_data), and map each slide to a suitable layout from the PowerPoint template txt.

    You are given:
    1. slide_data — a list of slides, each containing keys like title, bullet_points, and image, along with two additional keys to be filled:
    - layout_index: (currently empty or placeholder)
    - placeholders: a dict to be filled with correct placeholder indices used for "title", "content", and optionally "image".
    {ppt_json_dict}
    


    2. layout_info — extracted from a PowerPoint template. It contains all available layout_index, `layout_name`, and each placeholder's idx and type (like title, content, picture).
    {ppt_template_layout_info}
    Your task:
    - Match each slide with the  most appropriate layout using the content structure:
    - If a slide page  is about title page, match a layout with a title-only placeholder .
    - If it has title + bullet_points, find a layout with title + content or similar.
    - If it has an image, use a layout that includes an image or picture placeholder**.
    - For title + bullet_points + image, match with a layout that has title + content + image/picture.

    Rules:
    - Do not modify the slide's dicitonary existing content.
    - Just add the layout_index and placeholders fields in-place for each slide.
    - Use exact placeholder indices (`placeholder_format.idx`) from the `layout_info` — do not invent or assume numbers.
    - If no perfect layout is found, pick the closest match and fill as best as possible.


    STRICT SLIDE STRUCTURE FORMAT

    {{
        "title": "...",
        "bullet_points": [...],
        "image": "...",
        "layout_index": 2,
        "placeholder:{{
            "title":0,
            "content":1,
            "image":2
        }}

    
    }}

    """

    response = client.models.generate_content(
        model = "gemini-2.5-pro",
        contents=system_prompt_5,
    )

    final_ppt_text = response.text.strip()



    if final_ppt_text.startswith("```json"):
        final_ppt_text= re.sub(
            r"^```json\s*|\s*```$", "", final_ppt_text.strip(), flags=re.MULTILINE
        )


    try:
        final_ppt_dict = json.loads(final_ppt_text)
    except:
        print("Not valid json format")
        final_ppt_dict = ppt_json_dict



    with open("final_ppt_data.json","w",encoding="utf-8") as f:
        json.dump(final_ppt_dict,f,indent=4, ensure_ascii=False)


    print("\nPhase-7: Final Presenation Json generated\n")

    make_ppt_from_data(template_path,output_path)




if __name__ == "__main__":
    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    template_path = "../templates/15.pptx"
    print("............................\n")

    print(input_pdf)
    print(output_pptx)
    print(template_path)
    
    generate_ppt(input_path=input_pdf, output_path=output_pptx,template_path=template_path)
