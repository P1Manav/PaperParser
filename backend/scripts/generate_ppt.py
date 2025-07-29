import json
import sys
import os
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
import google.generativeai as genai
from dotenv import load_dotenv

# Load .env environment variables
load_dotenv()
sys.stdout.reconfigure(encoding='utf-8')  # Prevent encoding issues in console

# === Function to create PPT from JSON ===
def generate_ppt_from_template(json_content, template_path, output_path):
    print(f"🔸 Checking template path: {template_path}")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"❌ Template file not found at: {template_path}")

    print("🔸 Loading PowerPoint template...")
    prs = Presentation(template_path)

    try:
        slides_data = json.loads(json_content)
    except json.JSONDecodeError as e:
        raise ValueError(f"❌ Invalid JSON passed to generate_ppt_from_template: {e}")

    print("🔸 Populating slides...")
    for i, slide_data in enumerate(slides_data.get("slides", [])):
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

        if "content" in slide_data:
            try:
                content_box = slide.shapes.placeholders[1].text_frame
                content_box.clear()
                for point in slide_data["content"]:
                    content_box.add_paragraph().text = point.strip()
            except Exception as e:
                print(f"⚠️ Could not add bullet points: {e}")

    prs.save(output_path)
    print(f"✅ Presentation saved at: {output_path}")

# === Main Execution ===
if __name__ == "__main__":
    if len(sys.argv) < 2:
        raise ValueError("❌ No PDF file provided. Usage: python generate_ppt.py <pdf_path>")

    input_pdf = sys.argv[1]
    if not os.path.exists(input_pdf):
        raise FileNotFoundError(f"❌ PDF not found at: {input_pdf}")

    print("🔹 Extracting text from PDF...")
    loader = PyPDFLoader(input_pdf)
    text = " ".join([doc.page_content for doc in loader.load()])[:8000]  # Limit input

    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

    # First model: get summary
    print("🔹 Generating Summary...")
    summary_model = genai.GenerativeModel("gemini-2.0-flash-exp", generation_config={"response_mime_type": "text/plain"})
    summary_chat = summary_model.start_chat()
    summary_response = summary_chat.send_message(text).text.strip()

    if not summary_response:
        raise ValueError("❌ Empty AI summary response")

    print("🔹 AI Summary Response:")
    print(summary_response)

    # Second model: convert summary to JSON
    print("🔹 Generating JSON...")
    json_model = genai.GenerativeModel(
        model_name="gemini-2.0-flash-exp",
        generation_config={"response_mime_type": "application/json"},
        system_instruction="""Convert the research paper summary into a structured JSON presentation.
        Each slide should contain a title and up to 4-5 bullet points.
        Output format:
        {
            "presentation_title": "Research Paper Title",
            "slides": [
                {"type": "title", "title": "Title Slide", "subtitle": "Subtitle"},
                {"type": "content", "title": "Slide Title", "content": ["Point 1", "Point 2"]}
            ]
        }
        """
    )
    json_chat = json_model.start_chat()
    try:
        json_text = json_chat.send_message(summary_response).text.strip()
    except Exception as e:
        raise RuntimeError(f"❌ Gemini JSON generation failed: {e}")

    if not json_text:
        raise ValueError("❌ Empty AI JSON response")

    print("🔹 AI JSON Response:")
    print(json_text)

    # Validate JSON
    print("🔹 Validating JSON...")
    try:
        json.loads(json_text)
    except json.JSONDecodeError as e:
        print("❌ Bad JSON content:")
        print(json_text)
        raise ValueError(f"❌ Invalid JSON from Gemini: {e}")

    # Prepare paths
    template_path = os.path.join(os.path.dirname(__file__), "Template3-ppt.pptx")
    os.makedirs("outputs", exist_ok=True)
    output_pptx = os.path.join("outputs", "generated_presentation.pptx")

    print("🔹 Generating PowerPoint...")
    generate_ppt_from_template(json_text, template_path, output_pptx)

    # Optional: Clean up Gemini session
    # if hasattr(genai, "cleanup"):
    #     genai.cleanup()

    print(f"\n✅ PowerPoint generated successfully: {output_pptx}")
