import sys
import json
import os
import re
from pptx import Presentation
import fitz  # PyMuPDF
from langchain_community.document_loaders import PyMuPDFLoader
import google.generativeai as genai
from dotenv import load_dotenv
from pptx.util import Pt
from PIL import Image

load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")

if not GOOGLE_API_KEY:
    print("Error: GOOGLE_API_KEY or GEMINI_API_KEY environment variable not set")
    sys.exit(1)

genai.configure(api_key=GOOGLE_API_KEY)

if len(sys.argv) < 5:
    print("Usage: python ppt_gen.py <pdf_path> <output_pptx_path> <template_number> <length>")
    sys.exit(1)

pdf_path = sys.argv[1]
output_pptx_path = sys.argv[2]
template_number = sys.argv[3]
length_of_ppt = sys.argv[4]

if not os.path.exists(pdf_path):
    print(f"Error: PDF file not found: {pdf_path}")
    sys.exit(1)

# Template path inside scripts/templates
template_path = os.path.join(os.path.dirname(__file__), "templates", f"{template_number}.pptx")

if not os.path.exists(template_path):
    print(f"Error: Template file not found: {template_path}")
    sys.exit(1)

def make_ppt_from_data(template_path: str):
    """Create PowerPoint presentation from JSON data"""
    try:
        pre = Presentation(template_path)
        
        json_file = "final_ppt_data.json"
        if not os.path.exists(json_file):
            print(f"Error: {json_file} not found")
            return False
            
        with open(json_file, "r", encoding="utf-8") as f:
            slides_data = json.load(f)

        if not slides_data:
            print("Error: No slide data found")
            return False

        # First slide (title slide)
        title_slide_data = slides_data[0]
        layout_index = title_slide_data.get("layout_index", 0)
        placeholder = title_slide_data.get("placeholders", {})
        
        try:
            slide_layout = pre.slide_layouts[layout_index]
            slide = pre.slides.add_slide(slide_layout)
        except IndexError:
            print(f"Warning: Layout index {layout_index} not found, using layout 0")
            slide_layout = pre.slide_layouts[0]
            slide = pre.slides.add_slide(slide_layout)

        if "title" in placeholder:
            try:
                title_shape = slide.placeholders[placeholder["title"]]
                title_shape.text = title_slide_data["slide_title"]
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(32)
            except:
                print(" Title page main heading not found")

        if "content" in placeholder:
            try:
                content_shape = slide.placeholders[placeholder["content"]]
                content_shape.text_frame.text = ""
                for bullet in title_slide_data["bullet_points"]:
                    p = content_shape.text_frame.add_paragraph()
                    p.text = bullet
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(14)
            except:
                print(" Title page sub heading not found")

        # Remaining slides
        for slide_data in slides_data[1:]:
            layout_index = slide_data.get("layout_index", 0)
            placeholders = slide_data.get("placeholders", {})
            slide_layout = pre.slide_layouts[layout_index]
            slide = pre.slides.add_slide(slide_layout)

            if "title" in placeholders:
                try:
                    title_shape = slide.placeholders[placeholders["title"]]
                    title_shape.text = slide_data["slide_title"]
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(30)
                            run.font.bold = True
                except IndexError:
                    print(f" Title placeholder index {placeholders['title']} not found")

            if "content" in placeholders:
                try:
                    content_shape = slide.placeholders[placeholders["content"]]
                    content_shape.text_frame.text = ""
                    raw_bullets = slide_data.get("bullet_points", [])
                    cleaned_bullets = [bullet.replace("**", "").strip() for bullet in raw_bullets]
                    for bullet in cleaned_bullets:
                        p = content_shape.text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        for run in p.runs:
                            run.font.size = Pt(14)
                            run.font.bold = True
                except IndexError:
                    print(f" Content placeholder index {placeholders['content']} not found")

            if "image" in placeholders and slide_data["image_path"] != "null":
                try:
                    ph = slide.placeholders[placeholders["image"]]
                    img_path = os.path.join("images", slide_data["image_path"])
                    ph_left, ph_top, ph_width, ph_height = ph.left, ph.top, ph.width, ph.height
                    im = Image.open(img_path)
                    img_width, img_height = im.size
                    img_ratio = img_width / img_height
                    ph_ratio = ph_width / ph_height
                    if img_ratio > ph_ratio:
                        new_width = ph_width
                        new_height = ph_width / img_ratio
                    else:
                        new_height = ph_height
                        new_width = ph_height * img_ratio
                    left = ph_left + (ph_width - new_width) // 2
                    top = ph_top + (ph_height - new_height) // 2
                    slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)
                except Exception as e:
                    print(f" Image not inserted on slide '{slide_data['slide_title']}': {e}")

        pre.save(output_pptx_path)
        print(f" Final PPT Generated: {output_pptx_path}")
        return True
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        return False

def group_image_bboxes(image_bboxes, y_tol=80, x_gap_tol=100):
    # Sort images top to bottom, then left to right
    image_bboxes.sort(key=lambda b: (round(b.y0), b.x0))
    groups = []
    current_group = []

    for bbox in image_bboxes:
        if not current_group:
            current_group.append(bbox)
            continue

        last_bbox = current_group[-1]

        # Check if bbox is on same row (within y-tolerance)
        same_row = abs(bbox.y0 - last_bbox.y0) < y_tol

        # Check if horizontal gap is small enough
        horizontal_gap = bbox.x0 - last_bbox.x1
        close_enough = horizontal_gap < x_gap_tol

        if same_row and close_enough:
            current_group.append(bbox)
        else:
            groups.append(current_group)
            current_group = [bbox]

    if current_group:
        groups.append(current_group)

    return groups

def extract_combined_images_with_captions(file_path: str):
    pdf_file = fitz.open(file_path)
    os.makedirs("images", exist_ok=True)
    image_caption_map = {}
    image_id = 1  # to name images like image1.png, image2.png ...

    for page_index in range(len(pdf_file)):
        page = pdf_file.load_page(page_index)
        image_list = page.get_images(full=True)

        # Step 1: Collect image bounding boxes
        bboxes = []
        for img in image_list:
            image_name = img[7]  # get image name like 'Im1'
            try:
                bbox = page.get_image_bbox(image_name)
                bboxes.append(bbox)
            except Exception:
                continue

        # Step 2: Group images that are side-by-side
        groups = group_image_bboxes(bboxes)

        for group in groups:
            # Step 3: Merge bounding box of grouped images
            merged_bbox = fitz.Rect(
                min(b.x0 for b in group),
                min(b.y0 for b in group),
                max(b.x1 for b in group),
                max(b.y1 for b in group),
            )

            # Step 4: Extract combined image as pixmap
            zoom = 10
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, clip=merged_bbox)

            image_filename = f"image{image_id}.png"
            image_path = os.path.join("images", image_filename)
            pix.save(image_path)

            # Step 5: Extract caption just below the image group
            caption_area = fitz.Rect(
                merged_bbox.x0 - 300,
                merged_bbox.y1,
                merged_bbox.x1 + 300,
                merged_bbox.y1 + 400,  # Area below image to search for caption
            )
            caption_text = page.get_textbox(caption_area).strip()

            # Find "Figure X: ..." using regex
            figure_match = re.search(
                r"(?:Figure|Fig\.)\s?\d+\s?:\s?.+",
                caption_text,
                re.IGNORECASE,
            )

            figure_caption = (
                figure_match.group(0).strip()
                if figure_match
                else "No Figure caption found"
            )

            # Step 6: Save mapping
            image_caption_map[image_filename] = figure_caption

            image_id += 1

    # Step 7: Save all captions to JSON
    with open("images/image_captions.json", "w", encoding="utf-8") as f:
        json.dump(image_caption_map, f, indent=2, ensure_ascii=False)

def main():
    """Main execution function with proper error handling"""
    try:
        # Phase 1: Data extraction from PDF
        loader = PyMuPDFLoader(pdf_path)
        docs = loader.load()
        text = "".join(doc.page_content for doc in docs)
        print(" Phase-1: Text Extraction completed")

        # Phase 2: Text processing with Gemini
        type_pdf = "Research Paper"
        
        system_prompt_1 = f"""
You are a smart document structuring assistant. Your task is to take the complete extracted text from a PDF file and preprocess it in a way that it becomes well-organized, logically grouped, and ready for PowerPoint slide generation in the next phase.

Here are the inputs:

-----------------------------------
TYPE OF PDF: {type_pdf}
TARGET PPT LENGTH: {length_of_ppt} (short = 5-6 slides, medium = 7-8 slides, long = 10-12 slides)

FULL RAW TEXT EXTRACTED FROM PDF:
{text}
-----------------------------------

Your goals:

1. Understand and analyze the structure of the document based on its type:
   - For research papers, identify and label sections like: Abstract, Introduction, Problem Statement, Literature Review, Methodology, Results, Conclusion, Future Work, etc.
   - For standard documents (e.g. proposals, business plans), group and label based on logical themes or topics like Overview, Objectives, Key Concepts, Strategies, Implementation, Case Studies, etc.
   - If type is unknown or custom, use best judgment to extract and organize major topics or logical chunks.

2. Group the content into structured labeled sections:
   - Use headers like `## Section: Introduction`, `## Section: Results`, etc.
   - Each section should contain the full detailed text relevant to that label.
   - Preserve important technical or formal tone (especially for research-type PDFs)

3. Adjust the depth and number of sections based on the requested presentation length:
   - If short, limit to 5-6 major sections with concise, focused content.
   - If medium, allow 7-8 moderately detailed sections.
   - If long, include 10-12+ granular sections and sub-sections.

4. Preprocess the content:
   - Clean up repetitions or OCR errors.
   - Merge fragmented lines/sentences.
   - Retain any bullet points, formulas, or figure captions as-is.
   - If figure captions are found, label them clearly as `[Figure Caption: ...]`

5. Do NOT generate slides yet. Just return a well-organized, clean, and structured text document — ready for slide conversion in the next phase.

Return the output in clear Markdown or readable labeled format.

Now process the content accordingly.
"""

        model = genai.GenerativeModel("gemini-2.0-flash-exp")
        response = model.generate_content(system_prompt_1)
        pre_process_text = response.text
        print(" Phase-2: Text Processing completed")

        # Phase 3: Extract images
        extract_combined_images_with_captions(pdf_path)
        print(" Phase-3: Images extracted from PDF")

        # Phase 4: Extract figure captions
        system_prompt_2 = f"""
You are a figure caption extractor for research papers.

Below is the preprocessed and cleaned full text of a research paper. Your task is to accurately extract all figure captions, which follow standard academic formatting such as:

- "Figure 1: ..."  
- "Fig. 2: ..."  
- Sometimes long captions span multiple lines after the label.

INPUT TEXT:
{pre_process_text}

Your goals:

1. Extract every figure caption that starts with:
   - "Figure X:" (where X is a number)
   - OR "Fig. X:" (alternative abbreviation)

2. If a caption spans multiple lines, combine all the lines until the next paragraph or until the next "Figure"/"Fig." starts.

3. Return a clean ordered list of all figure captions found, in the format:
   - Caption Number (e.g., Figure 4)
   - Full Caption Text (include all lines that belong to it)

4. Do not modify or summarize captions. Just extract exactly what is written after the "Figure" or "Fig." label.

5. Skip any unrelated content.

Format your output like this:
Figure 1
Full caption text...

Figure 2
Another caption text.
"""

        response = model.generate_content(system_prompt_2)
        figure_captions = response.text
        print(" Phase-4: Figure captions extracted")

        # Phase 5: Process image captions
        image_captions_file = "images/image_captions.json"
        if os.path.exists(image_captions_file):
            with open(image_captions_file, "r", encoding="utf-8") as f:
                figure_captions_data = json.load(f)
        else:
            figure_captions_data = {}

        # Phase 6: Generate PPT data
        system_prompt_4 = f"""
You are a highly skilled research assistant and expert PowerPoint slide designer. Your task is to generate a structured JSON array representing a **professional, visually pleasing PowerPoint presentation** for a research paper. You will use the following two inputs

1. **preprocess_text**: A clean, linearized version of the full research paper. It contains all major sections — Abstract, Introduction, Methodology, Experiments, Results, Discussion, Figures,etc  and Conclusion — in plain text format.Below is complete pre_process_text

{pre_process_text}

2. **image_caption_dict**: A dictionary where the keys are image filenames (e.g., "image4.png") and the values are the corresponding cleaned figure captions. These images should only be placed on slides **where appropriate and contextually meaningful**.Given Below

{figure_captions_data}

YOUR OBJECTIVE

Generate a **parallel, clean, and consistent slide deck** structure in the form of a JSON array. Each object in the array represents a slide, including:

- Title
- Bullet points (if textual)
- Image (only when relevant)
- Professional styling attributes (font size, color)

The goal is to create a beautiful, reader-friendly research presentation that summarizes the key contributions of the paper section by section.

STRICT SLIDE STRUCTURE FORMAT

{
  "slide_title": "Introduction",

  "content_type": "bullet_points",  // or "image" or "mixed"
  "bullet_points": [
    "Each bullet should be concise and clear.",
    "Focus on key findings or motivations.",
    "Use simple language without jargon."
  ],

  "image_path": null  // or e.g., "image4.png" only if the figure caption is highly relevant
}
"""

        response = model.generate_content(system_prompt_4)
        ppt_json_text = response.text.strip()

        if ppt_json_text.startswith("\`\`\`json"):
            ppt_json_text = re.sub(
                r"^\`\`\`json\s*|\s*\`\`\`$", "", ppt_json_text.strip(), flags=re.MULTILINE
            )

        try:
            ppt_json_dict = json.loads(ppt_json_text)
        except:
            print("Not valid json format")
            ppt_json_dict = {}

        with open("ppt_data.json", "w", encoding="utf-8") as f:
            json.dump(ppt_json_dict, f, indent=4, ensure_ascii=False)

        print(" Phase-6: PPT data generated")

        # Phase 7: Alignment of the PPT
        pre = Presentation(template_path)
        ppt_template_layout_info = ""
        for i, layout in enumerate(pre.slide_layouts):
            ppt_template_layout_info += f"Layout {i} Name: {layout.name}\n"
            for j, placeholder in enumerate(layout.placeholders):
                ppt_template_layout_info += f"Placeholder {j} - idx: {placeholder.placeholder_format.idx}, type: {placeholder.placeholder_format.type}\n"

        system_prompt_5 = f"""
You're an AI assistant helping to generate a PowerPoint JSON structure. Use the provided dictionary (slide_data), and map each slide to a suitable layout from the PowerPoint template txt.

You are given:
1. slide_data — a list of slides, each containing keys like title, bullet_points, and image, along with two additional keys to be filled:
   - layout_index: (currently empty or placeholder)
   - placeholders: a dict to be filled with correct placeholder indices used for "title", "content", and optionally "image".
{ppt_json_dict}
   

2. layout_info — extracted from a PowerPoint template. It contains all available layout_index, `layout_name`, and each placeholder's idx and type (like title, content, picture).
{ppt_template_layout_info}
Your task:
- Match each slide with the  most appropriate layout using the content structure:
  - If a slide page  is about title page, match a layout with a title-only placeholder .
  - If it has title + bullet_points, find a layout with title + content or similar.
  - If it has an image, use a layout that includes an image or picture placeholder**.
  - For title + bullet_points + image, match with a layout that has title + content + image/picture.

Rules:
- Do not modify the slide's dicitonary existing content.
- Just add the layout_index and placeholders fields in-place for each slide.
- Use exact placeholder indices (`placeholder_format.idx`) from the `layout_info` — do not invent or assume numbers.
- If no perfect layout is found, pick the closest match and fill as best as possible.
"""

        response = model.generate_content(system_prompt_5)
        final_ppt_text = response.text.strip()

        if final_ppt_text.startswith("\`\`\`json"):
            final_ppt_text = re.sub(
                r"^\`\`\`json\s*|\s*\`\`\`$", "", final_ppt_text.strip(), flags=re.MULTILINE
            )

        try:
            final_ppt_dict = json.loads(final_ppt_text)
        except:
            print("Not valid json format")
            final_ppt_dict = ppt_json_dict

        with open("final_ppt_data.json", "w", encoding="utf-8") as f:
            json.dump(final_ppt_dict, f, indent=4, ensure_ascii=False)

        print(" Phase-7: Final Presentation JSON generated")

        # Final step: Create presentation
        success = make_ppt_from_data(template_path)
        if success:
            print(" Presentation generation completed successfully")
        else:
            print(" Presentation generation failed")
            sys.exit(1)
            
    except Exception as e:
        print(f" Error in main execution: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
