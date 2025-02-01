import json
import sys
import os
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
import google.generativeai as genai

# ✅ Ensure correct encoding to avoid Unicode errors
sys.stdout.reconfigure(encoding='utf-8')


# ✅ Function to create PPT from JSON
def generate_ppt_from_template(json_content, template_path, output_path):
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"❌ Template file not found at: {template_path}")

    prs = Presentation(template_path)  # Load the template
    
    try:
        slides_data = json.loads(json_content)  # Parse JSON content
    except json.JSONDecodeError as e:
        raise ValueError(f"❌ Invalid JSON: {e}\nReceived JSON:\n{json_content}")

    for i, slide_data in enumerate(slides_data["slides"]):
        slide_layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]

        if "content" in slide_data:
            text_box = slide.shapes.placeholders[1].text_frame
            text_box.clear()  # Remove default placeholder text
            for point in slide_data["content"]:
                p = text_box.add_paragraph()
                p.text = point.strip()

    prs.save(output_path)
    print(f"✅ Presentation saved at: {output_path}")
    sys.exit()

# ✅ Check if PDF path is provided
if len(sys.argv) < 2:
    raise ValueError("❌ No PDF file provided. Pass the PDF file path as an argument.")

input_pdf_path = sys.argv[1]

# ✅ Validate PDF file exists
if not os.path.exists(input_pdf_path):
    raise FileNotFoundError(f"❌ PDF file not found at: {input_pdf_path}")

# ✅ Load and extract text from the PDF
print("\n🔹 Extracting text from PDF...")
loader = PyPDFLoader(input_pdf_path)
docs = loader.load()
text = " ".join([doc.page_content for doc in docs])

# ✅ Limit text to prevent API overload
text = text[:8000]  # Adjust as needed

# ✅ Configure API key
api_key = "AIzaSyB7EOjDOS8d1-24aFdPWbGwjd3fIynADhs"  # Replace with actual key
genai.configure(api_key=api_key)

# ✅ Create AI model for summary
model = genai.GenerativeModel(
    model_name="gemini-2.0-flash-exp",
    generation_config={"response_mime_type": "text/plain"}
)

# ✅ Get summary from AI
print("\n🔹 Generating Summary...")
chat_session = model.start_chat()
response = chat_session.send_message(text)

if not response.text.strip():
    raise ValueError("❌ AI summary response is empty. Check API key or request format.")

print("\n🔹 AI Summary Response:")
print(response.text)  # Debugging

# ✅ Create AI model for structured JSON output
model_json = genai.GenerativeModel(
    model_name="gemini-2.0-flash-exp",
    generation_config={"response_mime_type": "application/json"},
    system_instruction="""Convert the research paper summary into a structured JSON presentation.
    Each slide should contain a title and up to 4-5 bullet points.
    Ensure the content is technical, factual, and structured correctly.
    Output should follow this format:
    {
        "presentation_title": "Research Paper Title",
        "slides": [
            {"type": "title", "title": "Research Paper Title", "subtitle": "Main Theme"},
            {"type": "content", "title": "Section Title", "content": ["Point 1", "Point 2"]} 
        ]
    }
    """
)

# ✅ Get JSON structure from AI
print("\n🔹 Generating JSON...")
chat_session_json = model_json.start_chat()
response_json = chat_session_json.send_message(response.text)

if not response_json.text.strip():
    raise ValueError("❌ AI JSON response is empty. Check API key or request format.")

print("\n🔹 AI JSON Response:")
print(response_json.text)  # Debugging

json_response_text = response_json.text.strip()

try:
    parsed_json = json.loads(json_response_text)  # Ensure valid JSON
except json.JSONDecodeError as e:
    raise ValueError(f"❌ AI returned invalid JSON: {e}\nReceived JSON:\n{json_response_text}")

# ✅ Define paths
template_selected = os.path.join(os.path.dirname(__file__), "Template3-ppt.pptx")
# Define the relative path to the folder
output_dir = os.path.join("outputs")

# Ensure the directory exists
os.makedirs(output_dir, exist_ok=True)

# Define the full path for the output file
output_pptx = os.path.join(output_dir, "generated_presentation.pptx")
# ✅ Generate PPT
generate_ppt_from_template(json_response_text, template_selected, output_pptx)

# ✅ Ensure proper cleanup after the process
genai.cleanup()  # If cleanup function exists to close API connections.

print(f"\n✅ PowerPoint generated successfully: {output_pptx}")

# Exit the program after generating the PowerPoint
  # Terminates the script after successful generation
